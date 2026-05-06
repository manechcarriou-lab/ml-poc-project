"""Generate the slide deck `deliverables/process_overview.pptx`.

The deck embeds images from `plots/` (run `scripts/generate_plots.py` first).

Run with:  python scripts/build_slides.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "deliverables" / "process_overview.pptx"
PLOTS = ROOT / "plots"

PRIMARY = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0xE6, 0x7E, 0x22)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)
DARK = RGBColor(0x2C, 0x3E, 0x50)
GREY = RGBColor(0x7F, 0x8C, 0x8D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS = RGBColor(0x27, 0xAE, 0x60)
SUBTITLE = RGBColor(0xCF, 0xDC, 0xE7)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))


def _add_text(box, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, name="Calibri"):
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = ""
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    return tf


def _add_bullets(box, items, size=15, color=DARK, name="Calibri", spacing=4):
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = ""
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = name
        p.space_after = Pt(spacing)


def _add_band(slide, top_in, height_in, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(top_in), Inches(13.333), Inches(height_in)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _new_slide(prs, title, subtitle=None, page=None, total=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_band(slide, 0, 0.85, PRIMARY)
    box = _add_textbox(slide, 0.5, 0.16, 12.3, 0.55)
    _add_text(box, title, size=24, bold=True, color=WHITE)
    if subtitle:
        sub = _add_textbox(slide, 0.5, 0.62, 12.3, 0.3)
        _add_text(sub, subtitle, size=12, color=SUBTITLE)
    if page is not None and total is not None:
        page_box = _add_textbox(slide, 11.7, 0.22, 1.3, 0.45)
        _add_text(page_box, f"{page} / {total}", size=11, color=WHITE, align=PP_ALIGN.RIGHT)
    return slide


def _add_table(slide, left_in, top_in, width_in, height_in, headers, rows, header_size=12, body_size=11):
    rows_n = len(rows) + 1
    cols_n = len(headers)
    shape = slide.shapes.add_table(
        rows_n, cols_n,
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in),
    )
    table = shape.table
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        run.font.bold = True
        run.font.size = Pt(header_size)
        run.font.color.rgb = WHITE
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(body_size)
            run.font.color.rgb = DARK
            if c == 0:
                run.font.bold = True
    return table


def _add_image(slide, image_path, left_in, top_in, width_in=None, height_in=None):
    if not Path(image_path).exists():
        return None
    kwargs = {}
    if width_in is not None:
        kwargs["width"] = Inches(width_in)
    if height_in is not None:
        kwargs["height"] = Inches(height_in)
    return slide.shapes.add_picture(str(image_path), Inches(left_in), Inches(top_in), **kwargs)


def _caption(slide, text, left_in, top_in, width_in, color=GREY, size=11):
    box = _add_textbox(slide, left_in, top_in, width_in, 0.4)
    _add_text(box, text, size=size, color=color, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    TOTAL = 16

    # 1. Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_band(slide, 0, 7.5, PRIMARY)
    _add_text(_add_textbox(slide, 0.8, 1.7, 11.7, 1.4),
              "Prédire la conversion d'un visiteur e-commerce",
              size=42, bold=True, color=WHITE)
    _add_text(_add_textbox(slide, 0.8, 3.1, 11.7, 0.7),
              "Du repo vide au modèle XGBoost — choix techniques et justifications",
              size=20, color=SUBTITLE)
    _add_bullets(_add_textbox(slide, 0.8, 5.5, 11.7, 1.3),
                 ["Manech Carriou — Albert School",
                  "Projet ML PoC — fork du template thom1100/ml-poc-project",
                  "Repo : github.com/manechcarriou-lab/ml-poc-project"],
                 size=14, color=SUBTITLE, spacing=6)

    # 2. Executive summary
    s = _new_slide(prs, "TL;DR", "Ce qu'il faut retenir en 30 secondes", 2, TOTAL)
    _add_table(s, 0.5, 1.2, 12.3, 1.6,
               ["Modèle retenu", "Encoder", "Threshold", "F1 test", "ROC-AUC", "Critère"],
               [["XGBoost", "Ordinal", "0.305", "0.6731", "0.9292", "✅ +12% vs cible 0.60"]],
               header_size=13, body_size=14)
    box = _add_textbox(s, 0.5, 3.0, 12.3, 0.5)
    _add_text(box, "Pipeline reproductible en 4 commandes — 3 niveaux d'optimisation", size=15, bold=True, color=ACCENT)
    box = _add_textbox(s, 0.5, 3.4, 12.3, 4)
    _add_bullets(box, [
        "Dataset public (UCI Online Shoppers) — 12 330 sessions, classes déséquilibrées 85/15",
        "Pipeline sklearn ColumnTransformer = zéro fuite test → train (validé par tests unitaires)",
        "1) Optuna (TPE) — 15 trials par modèle, espaces continus",
        "2) Encoder optimal par famille (XGBoost → Ordinal, +2 F1 pts)",
        "3) TunedThresholdClassifierCV — seuil de décision appris en CV sur le train",
        "MLflow tracking + Streamlit demo + 7 tests unitaires verts",
    ], size=14)

    # 3. Business
    s = _new_slide(prs, "Le problème business", "Pourquoi ce sujet", 3, TOTAL)
    _add_text(_add_textbox(s, 0.5, 1.1, 12.3, 1.0),
              "Peut-on prédire dès le début d'une session qu'un visiteur va acheter ?",
              size=20, bold=True, color=PRIMARY)
    _add_bullets(_add_textbox(s, 0.5, 2.3, 12.3, 4.5), [
        "Conversion moyenne e-commerce : 1 à 3 % — les sessions à fort potentiel sont rares",
        "Identifier ces sessions permet de cibler le retargeting / coupons / pop-ups",
        "Décision opérationnelle : déclencher (ou non) une action marketing pendant la session",
        "Stakeholder : équipe Growth / CRM d'un retailer en ligne",
        "Type de problème ML : classification binaire supervisée (Revenue ∈ {True, False})",
    ], size=16)

    # 4. Dataset
    s = _new_slide(prs, "Le dataset", "Online Shoppers Purchasing Intention (UCI)", 4, TOTAL)
    _add_table(s, 0.5, 1.1, 6.5, 3.2,
               ["Champ", "Valeur"],
               [["Source", "UCI ML Repository (CC BY 4.0)"],
                ["Volume", "12 330 sessions × 18 colonnes"],
                ["Features", "10 numériques + 7 catégorielles"],
                ["Cible", "Revenue (booléen)"],
                ["Balance", "84.5 % False / 15.5 % True"]])
    img_path = PLOTS / "eda_target_balance.png"
    _add_image(s, img_path, 7.4, 1.1, width_in=5.4)
    box = _add_textbox(s, 0.5, 4.6, 6.5, 0.5)
    _add_text(box, "Limites assumées", size=15, bold=True, color=ACCENT)
    _add_bullets(_add_textbox(s, 0.5, 5.0, 6.5, 2.2), [
        "Pas d'identifiant utilisateur → pas de parcours multi-session",
        "Pas d'année → pas de saisonnalité multi-annuelle",
        "Modalités catégorielles anonymisées (Region 1-9, OS 1-8)",
    ], size=13)
    _caption(s, "Distribution de la cible — fort déséquilibre", 7.4, 6.7, 5.4)

    # 5. Setup
    s = _new_slide(prs, "Setup technique", "Stack et reproductibilité", 5, TOTAL)
    _add_table(s, 0.5, 1.1, 12.3, 3.4,
               ["Composant", "Choix", "Pourquoi"],
               [["Versionning", "Git + GitHub", "Standard, fork du template du cours"],
                ["Auth GitHub", "Clé SSH ed25519", "Pas de token à gérer, sécurisé"],
                ["CLI GitHub", "gh", "Forker, push, PR depuis le terminal"],
                ["Python", "3.13 dans .venv", "Isolation des dépendances"],
                ["Tracking ML", "MLflow 3.11", "Persistance des essais + UI web"],
                ["Hyperparam search", "Optuna 4.8", "Recherche bayésienne (TPE)"]])
    _add_text(_add_textbox(s, 0.5, 4.8, 12.3, 0.4),
              "Reproductibilité — 4 commandes", size=15, bold=True, color=ACCENT)
    _add_bullets(_add_textbox(s, 0.5, 5.2, 12.3, 2.0), [
        "git clone git@github.com:manechcarriou-lab/ml-poc-project.git",
        "python -m venv .venv && .venv\\Scripts\\activate",
        "pip install -r requirements.txt",
        "python scripts/train.py --trials 15",
    ], size=13, name="Consolas")

    # 6. EDA — checks
    s = _new_slide(prs, "EDA — checks de qualité", "data_exploration.ipynb", 6, TOTAL)
    _add_table(s, 0.5, 1.1, 12.3, 3.8,
               ["Check", "Résultat", "Décision"],
               [["NaN globaux", "0 ligne, 0 valeur", "Pas d'imputation"],
                ["Features avec >5 % NaN", "aucune", "Pas de drop"],
                ["Outliers IQR", "PageValues, BounceRates, *_Duration", "log1p + scaling"],
                ["Drift", "SpecialDay seulement au début", "À surveiller"],
                ["Imbalance cible", "84.5 / 15.5", "stratify=y + class_weight=balanced"],
                ["Modalités rares", "VisitorType=Other (0.7 %)", "OneHotEncoder(handle_unknown=ignore)"]])
    _add_text(_add_textbox(s, 0.5, 5.2, 12.3, 0.5),
              "Insight clé : New_Visitor convertit 2× plus que Returning_Visitor — feature très discriminante.",
              size=14, color=ACCENT, bold=True)

    # 7. EDA — segments (image)
    s = _new_slide(prs, "EDA — taux de conversion par segment", "Comprendre où sont les acheteurs", 7, TOTAL)
    _add_image(s, PLOTS / "eda_conversion_by_segment.png", 0.7, 1.3, width_in=12.0)
    _caption(s, "Taux de conversion par mois (gauche) et par type de visiteur (droite). Ligne rouge = moyenne globale.",
             0.5, 5.7, 12.3)
    _add_bullets(_add_textbox(s, 0.5, 6.2, 12.3, 1.2), [
        "Pic Nov-Sep-Oct : effet Black Friday + rentrée",
        "New_Visitor 2× au-dessus de la moyenne — contre-intuitif mais signal fort",
    ], size=13)

    # 8. Feature engineering pipeline
    s = _new_slide(prs, "Feature engineering", "Pipeline anti-leakage par construction", 8, TOTAL)
    _add_text(_add_textbox(s, 0.5, 1.1, 12.3, 0.5),
              "src/features.py — un seul ColumnTransformer dans un Pipeline sklearn",
              size=15, color=PRIMARY, bold=True)
    _add_bullets(_add_textbox(s, 0.5, 1.7, 12.3, 2.6), [
        "Skewed numeric (durations, PageValues) → log1p puis StandardScaler",
        "Numeric standard → StandardScaler",
        "Catégoriels → OneHotEncoder(handle_unknown='ignore')",
        "Engineered binaires → passthrough",
    ], size=14)
    _add_text(_add_textbox(s, 0.5, 4.5, 12.3, 0.5),
              "Pourquoi un Pipeline ?", size=15, bold=True, color=ACCENT)
    _add_bullets(_add_textbox(s, 0.5, 4.9, 12.3, 2.2), [
        "fit() est appelé uniquement sur le train (et sur les folds train de la CV)",
        "transform() applique au test les statistiques apprises sur le train",
        "→ aucune fuite test → train possible. Garantie standard sklearn.",
    ], size=14)

    # 9. Features added
    s = _new_slide(prs, "Features ajoutées", "Transformations row-wise stateless (leakage-safe)", 9, TOTAL)
    _add_table(s, 0.5, 1.1, 12.3, 5.0,
               ["Feature", "Définition", "Intuition métier"],
               [["TotalPages", "Σ des compteurs de pages", "Volume global d'engagement"],
                ["TotalDuration", "Σ des durations", "Temps passé total"],
                ["AvgTimePerPage", "duration / pages", "Profondeur de lecture"],
                ["ProductRelatedRatio", "pages produit / total", "Focus produit"],
                ["HighPageValue", "PageValues > 0", "Flag session marchande"],
                ["IsHighBounce", "BounceRates > Q3", "Flag session zappée"],
                ["IsSpecialDay", "SpecialDay > 0", "Jour spécial"]],
               body_size=12)
    _add_text(_add_textbox(s, 0.5, 6.4, 12.3, 0.5),
              "Toutes ces transformations sont déterministes ligne par ligne — pas de fit, donc pas de leakage.",
              size=12, color=GREY)

    # 10. Models + tuning
    s = _new_slide(prs, "Optuna + MLflow", "scripts/train.py — la boucle d'expérimentation", 10, TOTAL)
    _add_text(_add_textbox(s, 0.4, 1.1, 6.0, 0.5),
              "Optuna — pourquoi", size=16, bold=True, color=PRIMARY)
    _add_bullets(_add_textbox(s, 0.4, 1.6, 6.0, 4.8), [
        "Recherche bayésienne (TPE) > Grid Search",
        "Échantillonne les bons coins, ignore les mauvais",
        "Espaces continus (learning_rate log-uniforme)",
        "15 trials Optuna ≈ 100 trials Grid Search",
    ], size=13)
    _add_text(_add_textbox(s, 6.9, 1.1, 6.0, 0.5),
              "MLflow — pourquoi", size=16, bold=True, color=PRIMARY)
    _add_bullets(_add_textbox(s, 6.9, 1.6, 6.0, 4.8), [
        "Trace persistante des essais (comparable J+15)",
        "Log auto : params + métriques + artefacts",
        "UI web : mlflow ui --backend-store-uri ./mlruns",
        "Artefact = le joblib du pipeline final",
    ], size=13)
    _add_text(_add_textbox(s, 0.5, 6.5, 12.3, 0.5),
              "→ 52 runs trackés dans cette étude (3 studies + 45 trials + 3 finals + extras)",
              size=12, color=ACCENT, align=PP_ALIGN.CENTER)

    # 11. Results table + ROC
    s = _new_slide(prs, "Résultats finaux", "Test set 2 466 sessions — pipeline avec threshold tuning intégré", 11, TOTAL)
    _add_table(s, 0.3, 1.1, 12.7, 2.8,
               ["Modèle", "Encoder", "Threshold", "F1", "Precision", "Recall", "ROC-AUC"],
               [["Logistic Regression", "OneHot", "0.244", "0.6426", "0.5786", "0.7225", "0.9137"],
                ["Random Forest", "OneHot", "0.360", "0.6683", "0.6403", "0.6990", "0.9204"],
                ["XGBoost (winner)", "Ordinal", "0.305", "0.6731", "0.6203", "0.7356", "0.9292"]],
               header_size=11, body_size=11)
    badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(4.2), Inches(6.3), Inches(1.0))
    badge.fill.solid(); badge.fill.fore_color.rgb = SUCCESS; badge.line.fill.background()
    tf = badge.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Critère atteint — F1 > 0.60 visé, 0.6731 obtenu (XGBoost) → +12 %"
    run.font.size = Pt(16); run.font.bold = True; run.font.color.rgb = WHITE
    _add_text(_add_textbox(s, 0.5, 5.5, 12.3, 0.6),
              "Lecture business : sur 100 acheteurs réels, on en attrape 74 (recall 0.74). On gagne +14 pts vs threshold 0.5.",
              size=13, color=GREY)

    # 12. Results — ROC + PR side by side
    s = _new_slide(prs, "Résultats — courbes ROC et Precision-Recall", "Comparaison des 3 modèles sur le test", 12, TOTAL)
    _add_image(s, PLOTS / "roc_curves.png", 0.5, 1.1, width_in=6.2)
    _add_image(s, PLOTS / "pr_curves.png", 6.8, 1.1, width_in=6.2)
    _caption(s, "ROC : XGBoost domine avec AUC=0.93. PR : précision >0.7 jusqu'à recall ~0.6 — modèle utilisable en pratique.",
             0.5, 6.4, 12.3, color=DARK, size=12)

    # 13. Confusion matrix XGBoost + Feature importance
    s = _new_slide(prs, "Le modèle XGBoost en détail", "Erreurs et features importantes", 13, TOTAL)
    _add_image(s, PLOTS / "confusion_matrix_xgboost.png", 0.5, 1.1, width_in=5.6)
    _add_image(s, PLOTS / "feature_importance_xgb.png", 6.5, 1.1, width_in=6.4)
    _caption(s, "Matrice de confusion (gauche) — TP / FP / FN / TN sur le test set",
             0.5, 6.4, 5.6, size=11)
    _caption(s, "Top 15 features par importance gain — PageValues domine largement",
             6.5, 6.4, 6.4, size=11)

    # 14. Threshold tuning
    s = _new_slide(prs, "Threshold tuning intégré au pipeline", "TunedThresholdClassifierCV — anti-leakage par construction", 14, TOTAL)
    _add_image(s, PLOTS / "threshold_tuning_xgb.png", 1.5, 1.1, width_in=10.3)
    _caption(s, "Seuil retenu pour XGBoost : 0.305 (appris en CV sur le train). F1 = 0.6731 vs 0.6562 à 0.5.",
             0.5, 5.7, 12.3, color=DARK, size=12)
    _add_bullets(_add_textbox(s, 0.5, 6.3, 12.3, 1.1), [
        "Chaque modèle wrappé dans TunedThresholdClassifierCV(scoring='f1', cv=5) — fit sur train uniquement",
        "Gain de +4 à +7 points de F1 selon le modèle, intégré au pipeline production",
    ], size=13)

    # 15. Streamlit demo
    s = _new_slide(prs, "Démo Streamlit", "src/app.py — comment je présenterais ça à un client", 15, TOTAL)
    _add_bullets(_add_textbox(s, 0.5, 1.2, 12.3, 5.5), [
        "Section 1 — Contexte business + KPIs (sessions, taux de conversion, métriques modèle)",
        "Section 2 — EDA interactive : class imbalance, conversion par segment (Month/VisitorType/...), top corrélations",
        "Section 3 — Comparaison des modèles : tableau + barchart par métrique + plots ROC/PR/CM",
        "Section 4 — Démo : sliders sur PageValues / BounceRates / VisitorType / Month → proba prédite + jauge",
        "Section 5 — Threshold playground : matrice de confusion live sur le test à seuil ajustable",
        "Lancement : python scripts/main.py  (ou)  streamlit run src/app.py",
    ], size=14)

    # 16. Takeaways + next steps
    s = _new_slide(prs, "À retenir + prochaines étapes", "Synthèse", 16, TOTAL)
    _add_text(_add_textbox(s, 0.5, 1.1, 12.3, 0.5),
              "À retenir", size=16, bold=True, color=PRIMARY)
    _add_bullets(_add_textbox(s, 0.5, 1.6, 12.3, 3.0), [
        "Pipeline sklearn ColumnTransformer = garantie zéro fuite test → train (validé par tests unitaires)",
        "Trois leviers d'optimisation cumulés : Optuna + encoder par famille + threshold tuning CV",
        "XGBoost final : F1 = 0.6731, ROC-AUC = 0.9292, recall = 0.7356 (74 % des acheteurs captés)",
        "MLflow trace 52+ runs et permet de comparer toutes les configurations a posteriori",
        "Tout est reproductible en 4 commandes — tests + Streamlit + plots regénérables à la demande",
    ], size=14)
    _add_text(_add_textbox(s, 0.5, 4.7, 12.3, 0.5),
              "Prochaines étapes", size=16, bold=True, color=ACCENT)
    _add_bullets(_add_textbox(s, 0.5, 5.2, 12.3, 2.0), [
        "Affiner XGBoost avec un budget Optuna élargi (50-100 trials) sur l'encoder Ordinal",
        "Ajouter PR-AUC + Brier score dans le panel de métriques (calibration des probas)",
        "Calibration de probabilités (CalibratedClassifierCV) si déploiement en scoring probabiliste",
        "Test A/B en prod : threshold 0.305 vs threshold métier (à définir avec l'équipe Growth)",
    ], size=14)
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.0), Inches(13.333), Inches(0.5))
    accent.fill.solid(); accent.fill.fore_color.rgb = ACCENT; accent.line.fill.background()
    tf = accent.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Merci — Manech Carriou — github.com/manechcarriou-lab/ml-poc-project"
    run.font.size = Pt(14); run.font.bold = True; run.font.color.rgb = WHITE

    OUT.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(OUT)
        target = OUT
    except PermissionError:
        target = OUT.with_name("process_overview_new.pptx")
        prs.save(target)
        print(f"Note: {OUT.name} was locked (PowerPoint open). Saved as {target.name} instead.")
    print(f"Wrote {target.relative_to(ROOT)}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
